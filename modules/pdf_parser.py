import pdfplumber
from pptx import Presentation

def extract_text_from_pdf(uploaded_file):
    text = ""
    with pdfplumber.open(uploaded_file) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text:
                text += f"\n--- Slide {i+1} ---\n{page_text}"
    return text.strip()

def extract_text_from_pptx(uploaded_file):
    text = ""
    prs = Presentation(uploaded_file)
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text.strip() + "\n"
        if slide_text:
            text += f"\n--- Slide {i+1} ---\n{slide_text}"
    return text.strip()